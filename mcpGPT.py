import streamlit as st
from datetime import datetime
import time
import os
import json
import openai
from typing import Dict, List, Any, Union
import base64
from io import StringIO
import importlib.util
import glob
import ast
import sys
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
import pptx
import textwrap
import inspect

# Default configuration
DEFAULT_CONFIG = {
    "api_type": "azure",
    "api_base": "https://your-endpoint.openai.azure.com/",
    "api_key": "your-api-key-here",
    "api_version": "2023-03-15-preview",
    "model": "gpt-4o-mini"
}

# Application state
if 'config' not in st.session_state:
    st.session_state.config = DEFAULT_CONFIG.copy()
if 'conversation' not in st.session_state:
    st.session_state.conversation = []
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = {}
if 'available_tools' not in st.session_state:
    st.session_state.available_tools = {}

def ensure_string_content(content: Any) -> str:
    """Ensure the content is a valid string, handling null/None and other types"""
    if content is None:
        return "[No content]"
    if isinstance(content, str):
        return content
    if isinstance(content, (int, float)):
        return str(content)
    if isinstance(content, (dict, list)):
        try:
            return json.dumps(content, ensure_ascii=False)
        except:
            return str(content)
    return str(content)

# Utility functions
def save_config():
    """Save configuration to file"""
    with open('config.json', 'w') as f:
        json.dump(st.session_state.config, f)

def load_config():
    """Load configuration from file"""
    try:
        with open('config.json', 'r') as f:
            st.session_state.config.update(json.load(f))
    except FileNotFoundError:
        pass

def init_openai():
    """Initialize OpenAI client"""
    openai.api_type = st.session_state.config['api_type']
    openai.api_base = st.session_state.config['api_base']
    openai.api_key = st.session_state.config['api_key']
    openai.api_version = st.session_state.config['api_version']

def convert_to_string(value: Any) -> str:
    """Convert any value to string safely"""
    if isinstance(value, (int, float)):
        return str(value)
    elif isinstance(value, (list, dict)):
        return json.dumps(value, ensure_ascii=False)
    elif isinstance(value, pd.DataFrame):
        return value.to_markdown()
    elif value is None:
        return "None"
    return str(value)

# File processing functions
def extract_text_from_pdf(file):
    """Extract text from PDF"""
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_excel(file):
    """Extract text from Excel file"""
    df = pd.read_excel(file)
    return df.to_markdown()

def extract_text_from_word(file):
    """Extract text from Word document"""
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_ppt(file):
    """Extract text from PowerPoint presentation"""
    prs = pptx.Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def process_uploaded_file(file):
    """Process uploaded file based on type"""
    file_ext = file.name.split('.')[-1].lower()
    
    if file_ext == 'pdf':
        return extract_text_from_pdf(file)
    elif file_ext in ['xlsx', 'xls']:
        return extract_text_from_excel(file)
    elif file_ext == 'docx':
        return extract_text_from_word(file)
    elif file_ext == 'pptx':
        return extract_text_from_ppt(file)
    elif file_ext == 'txt':
        return file.read().decode('utf-8')
    else:
        return f"File content {file.name} not extracted (unsupported format)"

# Tool management functions
def load_tools():
    """Load tools from tools/ directory"""
    tools_dir = 'tools'
    os.makedirs(tools_dir, exist_ok=True)
    
    st.session_state.available_tools = {}
    
    for tool_path in glob.glob(os.path.join(tools_dir, 'tool-*.py')):
        try:
            tool_name = os.path.splitext(os.path.basename(tool_path))[0].replace('tool-', '')
            spec = importlib.util.spec_from_file_location(tool_name, tool_path)
            mod = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(mod)
            
            st.session_state.available_tools[tool_name] = {
                'function': mod.function_call,
                'schema': getattr(mod, 'function_schema', {}),
                'description': getattr(mod, 'description', "No description available"),
                'code': inspect.getsource(mod)
            }
        except Exception as e:
            st.error(f"Error loading tool {tool_path}: {str(e)}")

def get_tools_schema():
    """Return tools schema for OpenAI"""
    return [
        {
            "name": name,
            "description": info.get('description', f"Execute {name} function"),
            "parameters": info['schema']
        } for name, info in st.session_state.available_tools.items()
    ]

def execute_tool(tool_name: str, arguments: Dict) -> Dict:
    """Execute a tool and return standardized response"""
    try:
        if tool_name not in st.session_state.available_tools:
            return {
                "success": False,
                "content": f"Tool {tool_name} not found",
                "error": "Tool not found"
            }
        
        tool_func = st.session_state.available_tools[tool_name]['function']
        result = tool_func(**arguments)
        
        return {
            "success": True,
            "content": convert_to_string(result),
            "raw_result": result
        }
    except Exception as e:
        return {
            "success": False,
            "content": f"Error executing tool {tool_name}: {str(e)}",
            "error": str(e)
        }

def show_tool_creation():
    """Display tool creation interface"""
    st.header("🛠 Create New Tool")
    
    with st.expander("Tool Creation Guide"):
        st.markdown("""
        **Required structure for a tool:**
        1. A `function_call` function that implements the logic
        2. A `function_schema` in JSON format
        3. (Optional) A `description` string
        
        **Example Code:**
        ```python
        # Tool schema
        function_schema = {
            "type": "object",
            "properties": {
                "param1": {"type": "string", "description": "Parameter description"}
            },
            "required": ["param1"]
        }
        
        # Tool description
        description = "This tool does something useful"
        
        # Main function
        def function_call(param1: str):
            \"\"\"Does something with param1\"\"\"
            return f"Result: {param1}"
        ```
        """)
    
    with st.form("new_tool_form"):
        tool_name = st.text_input("Tool name (no spaces)")
        tool_description = st.text_area("Tool description")
        
        # Schema editor
        schema_code = st.text_area(
            "Tool JSON Schema",
            value=textwrap.dedent('''\
            {
                "type": "object",
                "properties": {
                    "param1": {
                        "type": "string",
                        "description": "Parameter description"
                    }
                },
                "required": ["param1"]
            }'''),
            height=200
        )
        
        # Function editor
        function_code = st.text_area(
            "Function Code",
            value=textwrap.dedent('''\
            def function_call(param1: str):
                """Does something with param1"""
                return f"Result: {param1}"'''),
            height=300
        )
        
        if st.form_submit_button("Create Tool"):
            try:
                # Validate JSON schema
                schema = json.loads(schema_code)
                
                # Create file content
                tool_content = f'''# tool-{tool_name}.py
# Description: {tool_description}

function_schema = {schema}

{function_code}
'''
                # Save file
                tool_path = os.path.join('tools', f'tool-{tool_name}.py')
                with open(tool_path, 'w', encoding='utf-8') as f:
                    f.write(tool_content)
                
                # Reload tools
                load_tools()
                st.success(f"Tool '{tool_name}' created successfully!")
            except json.JSONDecodeError:
                st.error("Invalid JSON schema - check syntax")
            except Exception as e:
                st.error(f"Error creating tool: {str(e)}")

def show_tool_management():
    """Display tool management interface"""
    st.header("🔧 Tool Management")
    
    tab1, tab2 = st.tabs(["Create Tool", "Existing Tools"])
    
    with tab1:
        show_tool_creation()
    
    with tab2:
        if not st.session_state.available_tools:
            st.warning("No tools available")
        else:
            for tool_name, tool_info in st.session_state.available_tools.items():
                with st.expander(f"🛠 {tool_name}"):
                    st.markdown(f"**Description:** {tool_info.get('description', 'No description')}")
                    
                    st.markdown("**Schema:**")
                    st.json(tool_info['schema'])
                    
                    st.markdown("**Function Code:**")
                    st.code(tool_info['code'], language='python')
                    
                    if st.button(f"Delete {tool_name}", key=f"del_{tool_name}"):
                        try:
                            os.remove(os.path.join('tools', f'tool-{tool_name}.py'))
                            load_tools()
                            st.success(f"Tool {tool_name} deleted!")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Error deleting tool: {str(e)}")

def chat_with_llm(messages: List[Dict]) -> Dict:
    """Send messages to OpenAI API with content validation"""
    try:
        # Prepare messages with validated content
        validated_messages = []
        for msg in messages:
            validated_msg = msg.copy()
            validated_msg['content'] = ensure_string_content(msg.get('content', ''))
            validated_messages.append(validated_msg)
        
        tools = get_tools_schema()
        response = openai.ChatCompletion.create(
            engine=st.session_state.config['model'],
            messages=validated_messages,
            tools=[{"type": "function", "function": t} for t in tools] if tools else None,
            tool_choice="auto" if tools else None,
        )
        
        return response.choices[0].message
    except Exception as e:
        st.error(f"OpenAI error: {str(e)}")
        return None

# UI Pages
def show_config_page():
    """Display API configuration page"""
    st.title("🔧 API Configuration")
    
    with st.form("api_config"):
        st.session_state.config['api_type'] = st.selectbox(
            "API Type",
            ["azure", "openai"],
            index=0 if st.session_state.config['api_type'] == "azure" else 1
        )
        
        st.session_state.config['api_base'] = st.text_input(
            "API Endpoint",
            value=st.session_state.config['api_base']
        )
        
        st.session_state.config['api_key'] = st.text_input(
            "API Key",
            type="password",
            value=st.session_state.config['api_key']
        )
        
        st.session_state.config['api_version'] = st.text_input(
            "API Version",
            value=st.session_state.config['api_version']
        )
        
        st.session_state.config['model'] = st.text_input(
            "Model",
            value=st.session_state.config['model']
        )
        
        if st.form_submit_button("Save Configuration"):
            save_config()
            init_openai()
            st.success("Configuration saved!")

def show_chat_page():
    """Display main chat page"""
    st.title("💬 Smart Chat")
    
    # Sidebar for files and tools
    with st.sidebar:
        st.header("📁 Files")
        uploaded_files = st.file_uploader(
            "Upload files",
            type=['pdf', 'xlsx', 'xls', 'docx', 'pptx', 'txt'],
            accept_multiple_files=True
        )
        
        for file in uploaded_files:
            if file.name not in st.session_state.uploaded_files:
                content = process_uploaded_file(file)
                st.session_state.uploaded_files[file.name] = content
                st.success(f"File {file.name} processed!")
        
        st.header("🛠 Tools")
        if st.button("Reload Tools"):
            load_tools()
            st.success("Tools reloaded!")
        
        if st.button("Manage Tools"):
            st.session_state.current_page = "Tool Management"
            st.rerun()
    
    # Display conversation
    for msg in st.session_state.conversation:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])
            if msg.get("timestamp"):
                st.caption(f"At {msg['timestamp']}")
    
    # Handle new messages
    if prompt := st.chat_input("Your message..."):
        now = datetime.now().strftime("%H:%M:%S")
        
        # Add user message with validated content
        user_msg = {
            "role": "user", 
            "content": ensure_string_content(prompt), 
            "timestamp": now
        }
        st.session_state.conversation.append(user_msg)
       
        with st.chat_message("user"):
            st.write(prompt)
            st.caption(f"At {now}")
        
        with st.spinner("Thinking..."):
            start_time = time.time()
            
            # Prepare context with uploaded files
            context = []
            if st.session_state.uploaded_files:
                context.append({
                    "role": "system",
                    "content": "Attached files:\n" + "\n\n".join(
                        f"=== {name} ===\n{content}" 
                        for name, content in st.session_state.uploaded_files.items()
                    )
                })
            
            # Add conversation history
            messages = context + [
                {"role": msg["role"], "content": msg["content"]} 
                for msg in st.session_state.conversation
                if msg["role"] in ["user", "assistant", "system"]
            ]
            
            # First LLM call
            response = chat_with_llm(messages)
            
            if response:
                # Handle tool calls
                if hasattr(response, 'tool_calls') and response.tool_calls:
                    # Execute tools
                    tool_responses = []
                    for call in response.tool_calls:
                        tool_name = call.function.name
                        args = json.loads(call.function.arguments)
                        
                        tool_result = execute_tool(tool_name, args)
                        
                        tool_responses.append({
                            "role": "tool",
                            "content": tool_result['content'],
                            "name": tool_name,
                            "tool_call_id": call.id
                        })
                    
                    # Add tool responses
                    messages.append({
                        "role": response.role,
                        "content": response.content,
                        "tool_calls": response.tool_calls
                    })
                    messages.extend(tool_responses)
                    
                    # Second call with tool results
                    final_response = chat_with_llm(messages)
                    
                    if final_response:
                        assistant_msg = {
                            "role": "assistant",
                            "content": final_response.content,
                            "timestamp": datetime.now().strftime("%H:%M:%S"),
                            "tools_used": [call.function.name for call in response.tool_calls]
                        }
                    else:
                        assistant_msg = {
                            "role": "assistant",
                            "content": "Error getting final response",
                            "timestamp": datetime.now().strftime("%H:%M:%S")
                        }
                else:
                    # Simple response without tools
                    assistant_msg = {
                        "role": "assistant",
                        "content": response.content,
                        "timestamp": datetime.now().strftime("%H:%M:%S")
                    }
                
                # Add to conversation
                st.session_state.conversation.append(assistant_msg)
                
                # Display response
                with st.chat_message("assistant"):
                    st.write(assistant_msg["content"])
                    st.caption(f"Response in {time.time()-start_time:.2f}s at {assistant_msg['timestamp']}")
                    if "tools_used" in assistant_msg:
                        st.info(f"Tools used: {', '.join(assistant_msg['tools_used'])}")

# Main application
def main():
    """Main application flow"""
    load_config()
    init_openai()
    load_tools()
    
    if 'current_page' not in st.session_state:
        st.session_state.current_page = "Chat"
    
    # Navigation sidebar
    st.sidebar.title("Navigation")
    st.session_state.current_page = st.sidebar.radio(
        "Go to",
        ["Chat", "API Configuration", "Tool Management"],
        index=["Chat", "API Configuration", "Tool Management"].index(st.session_state.current_page)
    )
    
    # Display current page
    if st.session_state.current_page == "Chat":
        show_chat_page()
    elif st.session_state.current_page == "API Configuration":
        show_config_page()
    elif st.session_state.current_page == "Tool Management":
        show_tool_management()

if __name__ == "__main__":
    main()
