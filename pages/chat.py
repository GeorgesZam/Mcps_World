import streamlit as st
from datetime import datetime
import time
import json
import openai
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
import pptx
from main import ensure_string_content

# Fonctions de traitement de fichiers
def extract_text_from_pdf(file):
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_excel(file):
    df = pd.read_excel(file)
    return df.to_markdown()

def extract_text_from_word(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_ppt(file):
    prs = pptx.Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def process_uploaded_file(file):
    file_ext = file.name.split('.')[-1].lower()
    
    if file_ext == 'pdf':
        return extract_text_from_pdf(file)
    elif file_ext in ['xlsx', 'xls']:
        return extract_text_from_excel(file)
    elif file_ext == 'docx':
        return extract_text_from_word(file)
    elif file_ext == 'pptx':
        return extract_text_from_ppt(file)
    elif file_ext == 'txt':
        return file.read().decode('utf-8')
    else:
        return f"File content {file.name} not extracted (unsupported format)"

# Fonctions de chat
def chat_with_llm(messages: list):
    try:
        validated_messages = []
        for msg in messages:
            validated_msg = msg.copy()
            validated_msg['content'] = ensure_string_content(msg.get('content', ''))
            validated_messages.append(validated_msg)
        
        response = openai.ChatCompletion.create(
            engine=st.session_state.config['model'],
            messages=validated_messages
        )
        return response.choices[0].message
    except Exception as e:
        st.error(f"OpenAI error: {str(e)}")
        return None

# Interface de chat
st.title("💬 Smart Chat")

# Gestion des fichiers
with st.sidebar:
    st.header("📁 Files")
    uploaded_files = st.file_uploader(
        "Upload files",
        type=['pdf', 'xlsx', 'xls', 'docx', 'pptx', 'txt'],
        accept_multiple_files=True
    )
    
    for file in uploaded_files:
        if file.name not in st.session_state.uploaded_files:
            content = process_uploaded_file(file)
            st.session_state.uploaded_files[file.name] = content
            st.success(f"File {file.name} processed!")

# Affichage de la conversation
for msg in st.session_state.conversation:
    with st.chat_message(msg["role"]):
        st.write(msg["content"])
        if msg.get("timestamp"):
            st.caption(f"At {msg['timestamp']}")

# Saisie de l'utilisateur
if prompt := st.chat_input("Your message..."):
    now = datetime.now().strftime("%H:%M:%S")
    
    user_msg = {
        "role": "user", 
        "content": ensure_string_content(prompt), 
        "timestamp": now
    }
    st.session_state.conversation.append(user_msg)
   
    with st.chat_message("user"):
        st.write(prompt)
        st.caption(f"At {now}")
    
    with st.spinner("Thinking..."):
        start_time = time.time()
        
        context = []
        if st.session_state.uploaded_files:
            context.append({
                "role": "system",
                "content": "Attached files:\n" + "\n\n".join(
                    f"=== {name} ===\n{content}" 
                    for name, content in st.session_state.uploaded_files.items()
                )
            })
        
        messages = context + [
            {"role": msg["role"], "content": msg["content"]} 
            for msg in st.session_state.conversation
            if msg["role"] in ["user", "assistant", "system"]
        ]
        
        response = chat_with_llm(messages)
        
        if response:
            assistant_msg = {
                "role": "assistant",
                "content": response.content,
                "timestamp": datetime.now().strftime("%H:%M:%S")
            }
            
            st.session_state.conversation.append(assistant_msg)
            
            with st.chat_message("assistant"):
                st.write(assistant_msg["content"])
                st.caption(f"Response in {time.time()-start_time:.2f}s at {assistant_msg['timestamp']}")
